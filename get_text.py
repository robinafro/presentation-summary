from pptx import Presentation
from PyPDF2 import PdfReader
from spire.doc import *
from spire.doc.common import *
from docx import Document as docx_document
import os 

def get_text(file_path):
    extension = os.path.splitext(file_path)[1]

    if extension == ".pptx":
        prs = Presentation(file_path)
        text_arr = []
        text = ""
        for slide_number, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_arr.append(shape.text)
                    text += shape.text + " "

        return text_arr, text
    elif extension == ".pdf":
        text_arr = []
        text = ""
        with open(file_path, "rb") as f:
            reader = PdfReader(f)
            for page_number in range(len(reader.pages)):
                text_arr.append(reader.pages[page_number].extract_text())
                text += reader.pages[page_number].extract_text() + " "

        return text_arr, text
    elif extension == ".docx":
        docx = docx_document(file_path)
        
        text = "\n".join(paragraph.text for paragraph in docx.paragraphs)
        
        return [], text
    elif extension == ".doc":
        document = Document()
        document.LoadFromFile(file_path)

        text = ("\n".join(document.GetText().split("\n")[1:])).strip()
        
        return [], text
